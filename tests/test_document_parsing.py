"""Tests for document text extraction utilities."""

from pathlib import Path

from nanobot.utils.document import (
    SUPPORTED_EXTENSIONS,
    _is_text_extension,
    extract_text,
)


class TestSupportedExtensions:
    """Test the SUPPORTED_EXTENSIONS constant."""

    def test_supported_extensions_include_common_formats(self):
        """Test that common document formats are included."""
        # Document formats
        assert ".pdf" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".xlsx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS

        # Text formats
        assert ".txt" in SUPPORTED_EXTENSIONS
        assert ".md" in SUPPORTED_EXTENSIONS
        assert ".csv" in SUPPORTED_EXTENSIONS
        assert ".json" in SUPPORTED_EXTENSIONS
        assert ".yaml" in SUPPORTED_EXTENSIONS
        assert ".yml" in SUPPORTED_EXTENSIONS

        # Image formats
        assert ".png" in SUPPORTED_EXTENSIONS
        assert ".jpg" in SUPPORTED_EXTENSIONS
        assert ".jpeg" in SUPPORTED_EXTENSIONS


class TestExtractText:
    """Test the extract_text function."""

    def test_extract_text_unsupported_returns_none(self, tmp_path: Path):
        """Test that unsupported file types return None."""
        unsupported_file = tmp_path / "file.xyz"
        unsupported_file.write_text("content")

        result = extract_text(unsupported_file)
        assert result is None

    def test_extract_text_file_not_found(self, tmp_path: Path):
        """Test that non-existent files return error string."""
        missing_file = tmp_path / "nonexistent.txt"

        result = extract_text(missing_file)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_txt_file(self, tmp_path: Path):
        """Test extracting text from a .txt file."""
        txt_file = tmp_path / "test.txt"
        content = "Hello, world!\nThis is a test."
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert result == content

    def test_extract_text_txt_file_with_truncation(self, tmp_path: Path):
        """Test that large text files are truncated."""
        txt_file = tmp_path / "large.txt"
        # Create content larger than _MAX_TEXT_LENGTH
        content = "x" * 300_000
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert len(result) < 300_000
        assert "(truncated," in result
        assert "chars total)" in result

    def test_extract_text_md_file(self, tmp_path: Path):
        """Test extracting text from a .md file."""
        md_file = tmp_path / "test.md"
        content = "# Header\n\nSome markdown content."
        md_file.write_text(content, encoding="utf-8")

        result = extract_text(md_file)
        assert result == content

    def test_extract_text_csv_file(self, tmp_path: Path):
        """Test extracting text from a .csv file."""
        csv_file = tmp_path / "test.csv"
        content = "name,age\nAlice,30\nBob,25"
        csv_file.write_text(content, encoding="utf-8")

        result = extract_text(csv_file)
        assert result == content

    def test_extract_text_json_file(self, tmp_path: Path):
        """Test extracting text from a .json file."""
        json_file = tmp_path / "test.json"
        content = '{"key": "value", "number": 42}'
        json_file.write_text(content, encoding="utf-8")

        result = extract_text(json_file)
        assert result == content

    def test_extract_text_xlsx(self, tmp_path: Path):
        """Test extracting text from an .xlsx file."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["A2"] = "Alice"
        ws["B2"] = 30
        ws["A3"] = "Bob"
        ws["B3"] = 25

        # Add a second sheet
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Product"
        ws2["B1"] = "Price"
        ws2["A2"] = "Widget"
        ws2["B2"] = 9.99

        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        assert result is not None
        assert "--- Sheet: Sheet1 ---" in result
        assert "--- Sheet: Sheet2 ---" in result
        assert "Alice" in result
        assert "Bob" in result
        assert "Widget" in result
        assert "9.99" in result

    def test_extract_text_xlsx_empty_sheet(self, tmp_path: Path):
        """Test extracting text from an .xlsx file with empty sheets."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "empty.xlsx"
        wb = Workbook()
        # Clear the default sheet
        wb.remove(wb.active)
        # Add an empty sheet
        wb.create_sheet("EmptySheet")
        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        # Empty sheets should return empty string or header only
        assert result == "--- Sheet: EmptySheet ---" or result == ""

    def test_extract_text_docx(self, tmp_path: Path):
        """Test extracting text from a .docx file."""
        from docx import Document

        docx_file = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is paragraph one.")
        doc.add_paragraph("This is paragraph two.")
        doc.save(docx_file)

        result = extract_text(docx_file)
        assert result is not None
        assert "Test Document" in result
        assert "This is paragraph one." in result
        assert "This is paragraph two." in result

    def test_extract_text_docx_empty(self, tmp_path: Path):
        """Test extracting text from an empty .docx file."""
        from docx import Document

        docx_file = tmp_path / "empty.docx"
        doc = Document()
        doc.save(docx_file)

        result = extract_text(docx_file)
        assert result == ""

    def test_extract_text_pptx(self, tmp_path: Path):
        """Test extracting text from a .pptx file."""
        from pptx import Presentation

        pptx_file = tmp_path / "test.pptx"
        prs = Presentation()

        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        for shape in slide1.shapes:
            if hasattr(shape, "text"):
                shape.text = "First Slide Title"

        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = width = height = 1000000
        textbox = slide2.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Bullet point content"

        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "--- Slide 1 ---" in result
        assert "--- Slide 2 ---" in result
        # Text content may vary depending on PowerPoint layout defaults
        assert len(result) > 0

    def test_extract_text_pptx_table(self, tmp_path: Path):
        """Table cells should be extracted, not silently dropped."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(1)
        ).table
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "Bob"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Header A" in result
        assert "Header B" in result
        assert "Alice" in result
        assert "Bob" in result

    def test_extract_text_pptx_grouped_shapes(self, tmp_path: Path):
        """Text inside grouped shapes must be extracted recursively."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "grouped.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        group = slide.shapes.add_group_shape()
        inner = group.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        inner.text_frame.text = "Inside group"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Inside group" in result

    def test_extract_text_pdf_not_found(self, tmp_path: Path):
        """Test that missing PDF files return error string."""
        missing_pdf = tmp_path / "nonexistent.pdf"

        result = extract_text(missing_pdf)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_image_files(self, tmp_path: Path):
        """Test that image files return placeholder text."""
        # Create a minimal PNG file (1x1 pixel)
        png_file = tmp_path / "test.png"
        # Minimal valid PNG: 8-byte signature + IHDR + IDAT + IEND
        png_data = (
            b"\x89PNG\r\n\x1a\n"
            b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01"
            b"\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        png_file.write_bytes(png_data)

        result = extract_text(png_file)
        assert result is not None
        assert "[image:" in result
        assert "test.png" in result


class TestIsTextExtension:
    """Test the _is_text_extension helper."""

    def test_text_extensions_return_true(self):
        """Test that known text extensions return True."""
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".md") is True
        assert _is_text_extension(".csv") is True
        assert _is_text_extension(".json") is True
        assert _is_text_extension(".yaml") is True
        assert _is_text_extension(".yml") is True
        assert _is_text_extension(".xml") is True
        assert _is_text_extension(".html") is True
        assert _is_text_extension(".htm") is True

    def test_non_text_extensions_return_false(self):
        """Test that non-text extensions return False."""
        assert _is_text_extension(".pdf") is False
        assert _is_text_extension(".docx") is False
        assert _is_text_extension(".xlsx") is False
        assert _is_text_extension(".pptx") is False
        assert _is_text_extension(".png") is False
        assert _is_text_extension(".xyz") is False

    def test_case_sensitivity(self):
        """Test that _is_text_extension requires lowercase extension.

        Note: The main extract_text function handles case-insensitivity by
        converting extensions to lowercase before calling _is_text_extension.
        """
        # _is_text_extension itself is case-sensitive (lowercase only)
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".TXT") is False
        assert _is_text_extension(".pdf") is False
