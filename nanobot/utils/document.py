"""Document text extraction utilities for nanobot."""

import mimetypes
from pathlib import Path

from loguru import logger

from nanobot.utils.helpers import detect_image_mime


# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    # Document formats
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    # Text formats
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    # Image formats (for future OCR support)
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
}

_MAX_TEXT_LENGTH = 200_000


def extract_text(path: Path) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    if not isinstance(path, Path):
        path = Path(path)

    if not path.exists():
        return f"[error: file not found: {path}]"

    ext = path.suffix.lower()

    # Document formats -- each branch lazily imports its parser so that
    # startup does not pay the ~25 MB cost of loading openpyxl /
    # python-docx / python-pptx / pypdf up front (see issue #3422).
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[error: pypdf not installed]"
    try:
        reader = PdfReader(path)
        pages: list[str] = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"--- Page {i} ---\n{text}")
        return _truncate("\n\n".join(pages), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to extract PDF {}", path)
        return f"[error: failed to extract PDF: {e!s}]"


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return "[error: python-docx not installed]"
    try:
        doc = DocxDocument(path)
        paragraphs: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        return _truncate("\n\n".join(paragraphs), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to extract DOCX {}", path)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[error: openpyxl not installed]"
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            return _truncate("\n\n".join(sheets), _MAX_TEXT_LENGTH)
        finally:
            wb.close()
    except Exception as e:
        logger.exception("Failed to extract XLSX {}", path)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        return "[error: python-pptx not installed]"
    try:
        prs = PptxPresentation(path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            if slide_text:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        return _truncate("\n\n".join(slides), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to extract PPTX {}", path)
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Collect text from a PPTX shape, recursing into groups and tables.

    Groups have ``has_text_frame=False`` and must be walked via ``.shapes``;
    tables are GraphicFrame objects whose cell text lives under ``.table``.
    """
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def _extract_text_file(path: Path) -> str:
    """Extract text from a plain text file."""
    try:
        # Try UTF-8 first, then latin-1 fallback
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to read text file {}", path)
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    """Truncate text with a suffix indicating truncation."""
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }


# ---------------------------------------------------------------------------
# High-level helper: split media into images + extracted document text
# ---------------------------------------------------------------------------

_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB


def extract_documents(
    text: str,
    media_paths: list[str],
    *,
    max_file_size: int = _MAX_EXTRACT_FILE_SIZE,
) -> tuple[str, list[str]]:
    """Separate images from documents in *media_paths*.

    Documents (PDF, DOCX, XLSX, PPTX, plain-text, …) have their text
    extracted and appended to *text*.  Only image paths are kept in the
    returned list so that downstream layers only need to handle vision
    blocks.

    Files larger than *max_file_size* bytes are skipped with a warning
    to avoid unbounded memory / CPU usage.
    """
    image_paths: list[str] = []
    doc_texts: list[str] = []

    for path_str in media_paths:
        p = Path(path_str)
        if not p.is_file():
            continue

        try:
            size = p.stat().st_size
        except OSError:
            continue
        if size > max_file_size:
            logger.warning(
                "Skipping oversized file for extraction: {} ({:.1f} MB > {} MB limit)",
                p.name, size / (1024 * 1024), max_file_size // (1024 * 1024),
            )
            continue

        with open(p, "rb") as f:
            header = f.read(16)
        mime = detect_image_mime(header) or mimetypes.guess_type(path_str)[0]
        if mime and mime.startswith("image/"):
            image_paths.append(path_str)
        else:
            extracted = extract_text(p)
            if extracted and not extracted.startswith("[error:"):
                doc_texts.append(f"[File: {p.name}]\n{extracted}")

    if doc_texts:
        text = text + "\n\n" + "\n\n".join(doc_texts)

    return text, image_paths
